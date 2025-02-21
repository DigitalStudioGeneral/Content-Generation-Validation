#Experiment
import streamlit as st
import os
from dotenv import load_dotenv
import json
from typing import Callable
from multiprocessing import Pool
from concurrent.futures import ThreadPoolExecutor, as_completed
import openai

# from langchain_google_vertexai import VertexAI
from langchain_openai import AzureOpenAI

from langchain import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain.docstore.document import Document as Document2 
from langchain.schema.runnable import RunnableMap
from langchain.schema.output_parser import StrOutputParser

import vertexai # Change this
import pdf2image
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
import docx
import pptx
# from langchain_community.vectorstores import Chroma
from langchain_core.prompts import PromptTemplate
# from langchain_google_vertexai import VertexAIEmbeddings #Change this
from langchain_openai import AzureOpenAIEmbeddings


import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.placeholder import PlaceholderPicture
import sys
import fitz
import zipfile
import shutil
from os.path import join
from os import listdir, rmdir
from shutil import move
import pandas as pd
from docx.api import Document

import base64
from langchain_community.chat_models import AzureChatOpenAI # Change this
from langchain_core.messages import HumanMessage
import uuid
from langchain.retrievers.multi_vector import MultiVectorRetriever
from langchain.storage import InMemoryStore
from langchain_core.documents import Document as Document1

import io
from io import BytesIO
import re

from IPython.display import HTML, display
from langchain_core.runnables import RunnableLambda, RunnablePassthrough
from PIL import Image
from st_multimodal_chatinput import multimodal_chatinput
from fpdf import FPDF
from pathlib import Path
from streamlit_extras.stylable_container import stylable_container



load_dotenv()

llm = AzureChatOpenAI (
        azure_endpoint=os.getenv("AZURE_ENDPOINT"),
        openai_api_key=os.getenv("AZURE_OPENAI_API_KEY"),
        
        azure_deployment=os.getenv("CHAT_DEPLOYMENT"),
        openai_api_version=os.getenv("CHAT_VERSION"),
        model_name=os.getenv("CHAT_MODEL"),
        temperature=0,
        )

embed_model = AzureOpenAIEmbeddings(
        openai_api_key=os.getenv("AZURE_OPENAI_API_KEY"),
        azure_endpoint=os.getenv("AZURE_ENDPOINT"),

        model=os.getenv("EMBEDDING_MODEL"),
        azure_deployment=os.getenv("EMBEDDING_DEPLOYMENT"),
        openai_api_version=os.getenv("EMBEDDING_VERSION"),
        dimensions=768,
        )

def extract_text_from_pdf(pdf_path, stream='no'):
    """Extracts text from a PDF file."""
    if stream=='no':
        with open(pdf_path, "rb", encoding="utf-8") as f:
            reader = PdfReader(f)
            text = "".join(page.extract_text() for page in reader.pages)
    if stream=='yes':
        text = []
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            text.append(page.extract_text())     
    return text

def extract_text_from_word(word_path):
    """Extracts text from a Word document."""
    doc = docx.Document(word_path)
    text = "".join(paragraph.text for paragraph in doc.paragraphs)
    return text

def extract_text_from_ppt(ppt_path):
    """Extracts text from a PowerPoint presentation."""
    prs = Presentation(ppt_path)
    text = "".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    return text

def extract_img_from_pdf(pdf_path, image_root_folder):
    doc = fitz.open(pdf_path)
    image_folder = image_root_folder
    if not os.path.exists(image_folder):
        os.makedirs(image_folder)
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Identity, dpi=None,
                              colorspace=fitz.csRGB, clip=None, alpha=True, annots=True)
        pix.save(image_folder + "/" + os.path.split(pdf_path)[-1].split(".")[0] + "_%i.png" % page.number)

def extract_img_from_word(word_path, image_root_folder):
    archive = zipfile.ZipFile(word_path)
    image_folder = image_root_folder
    if not os.path.exists(image_folder):
        os.makedirs(image_folder)
    for file in archive.filelist:
        if file.filename.startswith('word/media/') and file.file_size > 30000:
            archive.extract(file, image_root_folder)
        else:
            return

    for filename in listdir(join(image_root_folder, 'word/media/')):
        move(join(image_root_folder, 'word/media/', filename), join(image_folder + "/", os.path.split(word_path)[-1].split(".")[0] + "_" + filename))
    rmdir(join(image_root_folder, 'word/media/'))
    rmdir(join(image_root_folder, 'word/'))

def write_image(shape, slide_idx, image_idx, ppt_path, image_root_folder):
    image_folder = image_root_folder
    if not os.path.exists(image_folder):
        os.makedirs(image_folder)
    image = shape.image
    image_bytes = image.blob
    image_filename = f'{image_folder}/{os.path.split(ppt_path)[-1].split(".")[0]}_slide{slide_idx}_image{image_idx:03d}.{image.ext}'
    image_idx += 1
    with open(image_filename, 'wb', encoding="utf-8") as f:
        f.write(image_bytes)
    return image_idx

def visitor(shape, slide_idx, image_idx, ppt_path, image_root_folder):
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if isinstance(shape, PlaceholderPicture):
            image_idx = write_image(shape, slide_idx, image_idx, ppt_path, image_root_folder)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            image_idx = visitor(s, slide_idx, image_idx, ppt_path, image_root_folder)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_idx = write_image(shape, slide_idx, image_idx, ppt_path, image_root_folder)
    return image_idx

def extract_img_from_ppt(ppt_path, image_root_folder):
    prs = Presentation(ppt_path)
    img_count = 0
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            img_count = visitor(shape, idx, img_count, ppt_path, image_root_folder)

def extract_table_from_word(word_path):
    document = docx.Document(word_path)
    tables = document.tables
    list_of_tables = []
    list_of_data = []

    for table in tables:
        data = []

        for i, row in enumerate(table.rows):
            row_data = []

            if i == 0:
                headers = [cell.text.strip() for cell in row.cells]
                continue

            for cell in row.cells:
                text = cell.text.strip()
                row_data.append(text)

            if not headers:
                row_str = "\n".join(row_data)
            else:
                row_str = "\n".join([headers[i] + ":" + row_data[i] for i in range(len(headers))])

            data.append("".join(row_str))

        list_of_data.append("\n".join(data))

    list_of_tables.append("\n".join(list_of_data))

    return list_of_tables

def extract_table_from_pdf(pdf_path):
    pdf_file = open(pdf_path, 'rb')
    pdf_reader = PdfReader(pdf_file)

    table_data = []
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text = page.extract_text()
        lines = text.split('\n')

        for line in lines:
            table_data.append(line)

    return table_data

def extract_table_from_ppt(ppt_path):
    prs = Presentation((ppt_path))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            row_count = len(tbl.rows)
            col_count = len(tbl.columns)
            for r in range(0, row_count):
                for c in range(0, col_count):
                    cell = tbl.cell(r, c)
                    paragraphs = cell.text_frame.paragraphs
                    for paragraph in paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
    return text_runs

def process_file(file_path, image_root_folder):
    text = ""
    table_data = []
    if file_path.endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
        extract_img_from_pdf(file_path, image_root_folder)
        table_data = extract_table_from_pdf(file_path)
    elif file_path.endswith(".docx") or file_path.endswith(".doc"):
        text = extract_text_from_word(file_path)
        extract_img_from_word(file_path, image_root_folder)
        table_data = extract_table_from_word(file_path)
    elif file_path.endswith(".pptx") or file_path.endswith(".ppt"):
        text = extract_text_from_ppt(file_path)
        extract_img_from_ppt(file_path, image_root_folder)
        table_data = extract_table_from_ppt(file_path)
    return text, table_data

def process_files(folder_path, save_img):
    image_root_folder = os.path.join(save_img, "images")
    if not os.path.exists(image_root_folder):
        os.makedirs(image_root_folder)

    file_paths = [os.path.join(folder_path, filename) for filename in os.listdir(folder_path)]

    with Pool() as pool:
        results = pool.starmap(process_file, [(file_path, image_root_folder) for file_path in file_paths])

    texts = [result[0] for result in results]
    table_data = [result[1] for result in results]

    return texts, table_data

def plt_img_base64(img_base64):
    """Disply base64 encoded string as image"""
    # Create an HTML img tag with the base64 string as the source
    image_html = f'<img src="data:image/jpeg;base64,{img_base64}" />'
    # Display the image by rendering the HTML
    display(HTML(image_html))


def looks_like_base64(sb):
    """Check if the string looks like base64"""
    return re.match("^[A-Za-z0-9+/]+[=]{0,2}$", sb) is not None


def is_image_data(b64data):
    """
    Check if the base64 data is an image by looking at the start of the data
    """
    image_signatures = {
        b"\xFF\xD8\xFF": "jpg",
        b"\x89\x50\x4E\x47\x0D\x0A\x1A\x0A": "png",
        b"\x47\x49\x46\x38": "gif",
        b"\x52\x49\x46\x46": "webp",
    }
    try:
        header = base64.b64decode(b64data)[:8]  # Decode and get the first 8 bytes
        for sig, format in image_signatures.items():
            if header.startswith(sig):
                return True
        return False
    except Exception:
        return False


def resize_base64_image(base64_string, size=(128, 128)):
    """
    Resize an image encoded as a Base64 string
    """
    # Decode the Base64 string
    img_data = base64.b64decode(base64_string)
    img = Image.open(io.BytesIO(img_data))

    # Resize the image
    resized_img = img.resize(size, Image.LANCZOS)

    # Save the resized image to a bytes buffer
    buffered = io.BytesIO()
    resized_img.save(buffered, format=img.format)

    # Encode the resized image to Base64
    return base64.b64encode(buffered.getvalue()).decode("utf-8")


def split_image_text_types(docs):
    """
    Split base64-encoded images and texts
    """
    b64_images = []
    texts = []
    for doc in docs:
        # Check if the document is of type Document and extract page_content if so
        if isinstance(doc, Document2): #
            doc = doc.page_content
        if looks_like_base64(doc) and is_image_data(doc):
            doc = resize_base64_image(doc, size=(1300, 600))
            b64_images.append(doc)
        else:
            texts.append(doc)
    if len(b64_images) > 0:
        return {"images": b64_images[:1], "texts": []}
    return {"images": b64_images, "texts": texts}


def img_prompt_func(data_dict):
    """
    Join the context into a single string
    """
    formatted_texts = "\n".join(data_dict["context"]["texts"])
    messages = []

    # Adding the text for analysis
    text_message = {
        "type": "text",
        "text": (
            "You are digital marketing specialist and you are responsible for either generating content or criticing content based on the user ask.\n"
            "Step 1 : You have to identify if task is either for generation of content or validation of content. \n"
            "For generation usually user use words like generate, write, create, build, etc. and for Validation user use words like validate, check, approve, justify, give a proof, etc. \n"
            "Based on the task, you either follow task step 2a or step 2b and step2b.1. Make sure you only follow one of them not both."
            "Step 2a. If task is validation of content, then being digital marketing specialist, provide insight on how much a user given information are following the guidelines given in the context."
            "Check all the guideline on that given topic in the context and let the user know in detail which guideline is being followed and which is not.\n"
            "Be crticial and specially find out the points that doesn't adhere to given guidelines and sugguest possible improvements.\n"
            "Always make sure to additional check that content is using the  tone of voice of Cognizant in naration.\n"
            "or \n"
            "Step 2b. If task is Generation of content, then Generate content but do remember being digital marketing specialist you are responsible to generate content,"
            "for the given question and content should adhere to guidelines on that given topic provided in context.\n"
            "Always make sure to use tone of voice of Cognizant in naration while generating the content.\n"
            "Step 2b.1. Now perform Step 2a. on the result of Step 2b."
            "Step 3 : Always provide a score for how well content adhere to guidelines out of 10.\n"
            "You will be given a mixed of text, tables, and image(s) usually of charts or graphs.\n"
            f"User-provided question: {data_dict['question']}\n\n"
            "Text and / or tables:\n"
            f"{formatted_texts}"

        ),
    }
    messages.append(text_message)
    # Adding image(s) to the messages if present
    if data_dict["context"]["images"]:
        for image in data_dict["context"]["images"]:
            image_message = {
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{image}"},
            }
            # print(image_message)
            messages.append(image_message)
    return [HumanMessage(content=messages)]

# Multi vector Retriever
# Generate summaries of text elements
def generate_text_summaries(texts, tables, summarize_texts=False):
    """
    Summarize text elements
    texts: List of str
    tables: List of str
    summarize_texts: Bool to summarize texts
    """


    # Prompt
    prompt_text = """You are an assistant tasked with summarizing tables and text for retrieval. \
    These tables and text are guidelines. \
    These summaries will be embedded and used to retrieve the raw text or table elements. \
    Give a concise summary of the table or text that is well optimized for retrieval. Table or text: {element} """
    prompt = PromptTemplate.from_template(prompt_text)
    # model = VertexAI(
      # temperature=0, model_name="gemini-1.0-pro", max_output_tokens=1024)#.with_fallbacks([empty_response])
    model = llm
    summarize_chain = {"element": lambda x: x} | prompt | model | StrOutputParser()

    # Initialize empty summaries
    text_summaries = []
    table_summaries = []

    # Apply to text if texts are provided and summarization is requested
    if texts and summarize_texts:
        text_summaries = summarize_chain.batch(texts, {"max_concurrency": 1})
    elif texts:
        text_summaries = texts

    # Apply to tables if tables are provided
    if tables:
        table_summaries = summarize_chain.batch(tables, {"max_concurrency": 1})

    return text_summaries, table_summaries

def encode_image(image_path):
    """Getting the base64 string"""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

def image_summarize(img_base64, prompt):
    """Make image summary"""
    # model = ChatVertexAI(model_name="gemini-1.0-pro-vision", max_output_tokens=1024)
    model = llm

    msg = model(
        [
            HumanMessage(
                content=[
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{img_base64}"},
                    },
                ]
            )
        ]
    )
    return msg.content

def process_image(img_path, prompt):
    base64_image = encode_image(img_path)
    summary = image_summarize(base64_image, prompt)
    return base64_image, summary

def generate_img_summaries(path):
    """
    Generate summaries and base64 encoded strings for images
    path: Path to list of .jpg files extracted by Unstructured
    """

    # Store base64 encoded images
    img_base64_list = []

    # Store image summaries
    image_summaries = []

    # Prompt
    prompt = """You are an assistant tasked with summarizing images for retrieval. \
    These summaries will be embedded and used to retrieve the raw image. \
    Give a concise summary of the image that is well optimized for retrieval."""

    # Apply to images
    with ThreadPoolExecutor() as executor:
        futures = []
        for img_file in sorted(os.listdir(path)):
            if img_file.endswith(tuple([".jpg", ".png", ".jpeg"])):
                img_path = os.path.join(path, img_file)
                future = executor.submit(process_image, img_path, prompt)
                futures.append(future)

        for future in as_completed(futures):
            base64_image, summary = future.result()
            img_base64_list.append(base64_image)
            image_summaries.append(summary)

    return img_base64_list, image_summaries

def create_multi_vector_retriever(
    vectorstore, text_summaries, texts, table_summaries, tables, image_summaries, images
):
    """
    Create retriever that indexes summaries, but returns raw images or texts
    """

    # Initialize the storage layer
    store = InMemoryStore()
    id_key = "doc_id"

    # Create the multi-vector retriever
    retriever = MultiVectorRetriever(
        vectorstore=vectorstore,
        docstore=store,
        id_key=id_key,
    )

    # Helper function to add documents to the vectorstore and docstore
    def add_documents(retriever, doc_summaries, doc_contents):
        doc_ids = [str(uuid.uuid4()) for _ in doc_contents]
        summary_docs = [
            Document1(page_content=s, metadata={id_key: doc_ids[i]})
            for i, s in enumerate(doc_summaries)
        ]
        retriever.vectorstore.add_documents(summary_docs)
        retriever.docstore.mset(list(zip(doc_ids, doc_contents)))

    # Add texts, tables, and images
    # Check that text_summaries is not empty before adding
    if text_summaries:
        add_documents(retriever, text_summaries, texts)
    # Check that table_summaries is not empty before adding
    if table_summaries:
        add_documents(retriever, table_summaries, tables)
    # Check that image_summaries is not empty before adding
    if image_summaries:
        add_documents(retriever, image_summaries, images)

    return retriever


@st.cache_data
def extraction_pipeline(text_list, table_list, save_img):

    # model = VertexAI(model_name="gemini-1.0-pro", temperature =0.9)
    model = llm
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    context = "\n".join(str(p) for p in text_list)
    texts = text_splitter.split_text(context)

    # Get text, table summaries
    text_summaries, table_summaries = generate_text_summaries(
        texts, table_list, summarize_texts=True
    )

    while '' in text_summaries:
        text_summaries.remove('')
    while '' in table_summaries:
        table_summaries.remove('')

    print("Text and Table data is ready!!")

    # Image summaries
    img_base64_list, image_summaries = generate_img_summaries(f"./{save_img}")

    print("Image data is ready!!")

    print("Data is ready!!")

    return text_summaries, table_summaries, texts, table_list, img_base64_list, image_summaries


def multi_modal_rag_chain(retriever):
    """
    Multi-modal RAG chain
    """

    # Multi-modal LLM
    # model = ChatVertexAI(
        # temperature=0, model_name="gemini-1.0-pro-vision", max_output_tokens=1024)#1024
    model = llm

    # RAG pipeline
    chain = (
        {
            "context": retriever | RunnableLambda(split_image_text_types),
            "question": RunnablePassthrough(),
        }
        | RunnableLambda(img_prompt_func)
        | model
        | StrOutputParser()
    )

    print("chain has been created", chain)

    return chain


def main(retriever_multi_vector_img, username, guidelines_options, selected_guidelines, selected_document):

    if "current_file_index" not in st.session_state:
        st.session_state.current_file_index = 0
            
    counter = 0
    # Check if the username is set
    if username and selected_guidelines and selected_document:
        
        if selected_document == "I want to generate and validate content according to brand guidelines":
   
            chain_multimodal_rag_vg = multi_modal_rag_chain(retriever_multi_vector_img)
            print(chain_multimodal_rag_vg)
            print(type(chain_multimodal_rag_vg))

            output_container = st.container()
            chat_history = st.session_state.get("chat_history", [])
            if not os.path.exists("Chat_history"):
                os.makedirs("Chat_history")
            if not os.path.exists("user_input_image"):
                os.makedirs("user_input_image")
                
            
            
            
            #relative
            with stylable_container(
                key="bottom_content",
                css_styles="""
                    {
                        position: fixed;
                        bottom: 5px;
                        background-color: #141414;
                        width: 100%;
                    }
                    """,
            ):
                st.write(f"Considering {selected_guidelines} Guidelines")
                st.write("You can provide text and images for validation and generation.")
                chatinput = multimodal_chatinput()
                
            coded_image = []

            if chatinput:
                uploaded_file = chatinput["images"]
                user_input = chatinput["text"]
                if uploaded_file:
                    prompt = """You are an assistant tasked with summarizing images for retrieval. \
                            These summaries will be embedded and used to retrieve the raw image. \
                            Give a concise summary of the image that is well optimized for retrieval."""
                    parts = uploaded_file[0].split(',')
                    base64_data = parts[1] if len(parts) > 1 else ''
                    new_image_summaries = image_summarize(base64_data, prompt)
                    user_input = f"{user_input}. Image information:{new_image_summaries}"
                    coded_image.append(base64_data)

                if user_input:
                    chat_history.append(("human", user_input))
                    if coded_image:
                        chat_history.append(("human image", base64_data))
                    # else:
                    #     chat_history.append(("human image", "")) 
                    response = chain_multimodal_rag_vg.invoke(user_input)
                    chat_history.append(("ai", response))
                    st.session_state["chat_history"] = chat_history
                    
                    
                # with st.container:
                pdf = FPDF()
                # compression is not yet supported in py3k version
                pdf.compress = False
                pdf.add_page()
                # Unicode is not yet supported in the py3k version; use windows-1252 standard font
                pdf.add_font('DejaVu', '', 'DejaVuSans.ttf', uni=True)
                pdf.set_font('DejaVu', '', 14)  
                pdf.ln(10)
                with output_container:
                    for i in range(1,len(chat_history)+1):
                        sender, message = chat_history[len(chat_history) - i]
                        if(sender != "human image"):
                            # st.markdown(f"**{sender}:** {message}")
                            st.chat_message(sender).write(message)
                        if(sender == "human image" and message != ""):
                            st.markdown(f"""<img src="data:png;base64,{message}" width='50' height='50' >""", True)
                        if(i < 3):
                            if(sender != "human image"):
                                pdf.write(5, f"**{sender}:** {message}")
                            else: 
                                if(len(coded_image) != 0 and i==2):
                                    image_bytes = base64.b64decode(coded_image[0])
                                    image = Image.open(BytesIO(image_bytes))
                                    image.save(f"user_input_image/{username}_{counter}.png", format="PNG")
                                    pdf.add_page()
                                    pdf.image(f"user_input_image/{username}_{counter}.png", 10, 10)
                                    counter += 1
                    pdf.output(f'Chat_history/Chat_history_{username}.pdf', 'F')
                    print("created")
                
        elif selected_document == "I want to validate a document content according to brand guidelines":
            chain_multimodal_rag_vg = multi_modal_rag_chain(retriever_multi_vector_img)
            output_doc_container = st.container()
            st.session_state["chat_history"] = []
            chat_history = st.session_state.get("chat_history", [])
            if not os.path.exists("pdf_by_user"):
                os.makedirs("pdf_by_user")
            
            if "query" not in st.session_state:
                st.session_state["query"] = ""
            
            if 'clicked' not in st.session_state:
                st.session_state.submit = False
            # with st.container():
            with stylable_container(
                key="bottom_content_document",
                css_styles="""
                    {
                        position: fixed;
                        bottom: 5px;
                        background-color: #141414;
                        width: 100%
                        
                        }
                        
                    [data-testid='stFileUploader'] {
                            width: max-content;
                            }
                    """,
            ):
                col1, col2 = st.columns([2,4]) 
                with col1:
                    st.write(f"Considering {selected_guidelines} Guidelines")
                # st.write("You can provide pdf document for validation.")
                    query = st.text_input("Anything specific you want to validate in the pdf?(optional)", key='doc_val', value=st.session_state["query"])
                with col2:
                    file = st.file_uploader("You can provide pdf document for validation.", "pdf")#Please choose a pdf file
                doc_submit = st.button("Submit")#, on_click= st.session_state.submit = True)

                
            
            if doc_submit:
            
                if file != None:
                    chat_history.append(("human", file.name))
                    text = extract_text_from_pdf(file, stream='yes')
                    save_folder = 'pdf_by_user'
                    save_path = Path(save_folder, file.name)
                    with open(save_path, mode='wb') as w:
                        w.write(file.getvalue())

                    if save_path.exists():
                        st.success(f'File {file.name} is successfully saved!')



                    extract_img_from_pdf(f"{save_folder}/{file.name}", file.name.split('.', 1)[0])
                    print("images extracted")
                    prompt = """You are an assistant tasked with summarizing images for retrieval. \
                                These summaries will be embedded and used to retrieve the raw image. \
                                Give a concise summary of the image that is well optimized for retrieval."""
                    count = 0
                    with st.spinner('Wait for it...'):
                        for img_file in sorted(os.listdir(file.name.split('.', 1)[0])):
                            if img_file.endswith(tuple([".jpg", ".png", ".jpeg"])):
                                img_path = os.path.join(file.name.split('.', 1)[0], img_file)
                                base64_image = encode_image(img_path)
                                image_summaries = image_summarize(base64_image, prompt)
                                user_input = f'Can you validate this information. Query : {query}, Textual information : {text[count]} , Image Information : {image_summaries}'
                                response = chain_multimodal_rag_vg.invoke(user_input)
                                count += 1
                                chat_history.append(("ai", f"**Page {count}**, {response}"))
                                st.session_state["chat_history"] = chat_history
                    st.success('Done!')
                    
                    with output_doc_container:
                        for sender, message in chat_history:
                            print(chat_history)
                            st.chat_message(sender).write(message)
                else:
                    st.warning("Please upload a document.")
                    

                        
    else:
        st.warning("Please enter your name and select both the fields to enable the app functionality.")

if __name__ == "__main__":
    st.markdown(
        """
    <style>
        .st-emotion-cache-4oy321 {
            flex-direction: row-reverse;
            text-align: right;
        }
    </style>
    """,
        unsafe_allow_html=True,
    )
    
    guidelines_options = ["Website", "ExternalComms", "SocialMedia", "Emails"]

    if "username" not in st.session_state:
        st.session_state["username"] = ""
        
    sibebar_submit =False
    
    with st.sidebar:
        st.write("Just Create")
        with st.container():
            username = st.text_input("Enter Your Name", key='username', value=st.session_state["username"])
            document_options = ["I want to generate and validate content according to brand guidelines","I want to validate a document content according to brand guidelines"]
            selected_document  = st.radio("What do you want help with?", document_options, key="document_w_selection", index= None)
            guidelines_options = ["Website", "ExternalComms", "SocialMedia", "Emails"]
            selected_guidelines  = st.radio("Which guideline should we consider?", guidelines_options, key="guidelines_selection", index= None)
            
         
    if username and selected_guidelines and selected_document:
        vectordb = Chroma(persist_directory=f"./chroma/chroma_guidelines_{selected_guidelines}", embedding_function=embed_model)
        folder_path = f"Guidelines_subfolders/{selected_guidelines}"
        save_img = f"imgs_guidelines/imgs_guidelines_{selected_guidelines}"
        guidelines = f"Guidelines_subfolders/{selected_guidelines}"
        save_img_folder = f"imgs_guidelines/imgs_guidelines_{selected_guidelines}"

        if selected_guidelines == "Website":
            if os.path.exists(f"./chroma/chroma_guidelines_{selected_guidelines}"): 
                  retriever_multi_vector_img = vectordb.as_retriever()
            else:
                text_list_Website, table_list_Website = process_files(folder_path, save_img)
                text_summaries_Website, table_summaries_Website, texts_Website, table_list_Website, img_base64_list_Website, image_summaries_Website = extraction_pipeline(text_list_Website, table_list_Website, save_img_folder)


                retriever_multi_vector_img = create_multi_vector_retriever(vectordb, text_summaries_Website,texts_Website,
                                                                    table_summaries_Website,table_list_Website,
                                                                    image_summaries_Website,img_base64_list_Website,)
        elif selected_guidelines == "ExternalComms":

            if os.path.exists(f"./chroma/chroma_guidelines_{selected_guidelines}"):
                retriever_multi_vector_img = vectordb.as_retriever()
            else:
                text_list_ExternalComms, table_list_ExternalComms = process_files(folder_path, save_img)
                text_summaries_ExternalComms, table_summaries_ExternalComms, texts_ExternalComms, table_list_ExternalComms, img_base64_list_ExternalComms, image_summaries_ExternalComms = extraction_pipeline(text_list_ExternalComms, table_list_ExternalComms, save_img_folder)
                retriever_multi_vector_img = create_multi_vector_retriever(vectordb, text_summaries_ExternalComms,texts_ExternalComms,
                                                                    table_summaries_ExternalComms,table_list_ExternalComms,
                                                                    image_summaries_ExternalComms,img_base64_list_ExternalComms,)
        elif selected_guidelines == "SocialMedia":

            if os.path.exists(f"./chroma/chroma_guidelines_{selected_guidelines}"):
                retriever_multi_vector_img = vectordb.as_retriever()
            else:
                text_list_SocialMedia, table_list_SocialMedia = process_files(folder_path, save_img)
                text_summaries_SocialMedia, table_summaries_SocialMedia, texts_SocialMedia, table_list_SocialMedia, img_base64_list_SocialMedia, image_summaries_SocialMedia = extraction_pipeline(text_list_SocialMedia, table_list_SocialMedia, save_img_folder)
                retriever_multi_vector_img = create_multi_vector_retriever(vectordb, text_summaries_SocialMedia,texts_SocialMedia,
                                                                    table_summaries_SocialMedia,table_list_SocialMedia,
                                                                    image_summaries_SocialMedia,img_base64_list_SocialMedia,)
        elif selected_guidelines == "Emails":

            if os.path.exists(f"./chroma/chroma_guidelines_{selected_guidelines}"):
                retriever_multi_vector_img = vectordb.as_retriever()
            else: 
                text_list_Emails, table_list_Emails = process_files(folder_path, save_img)
                text_summaries_Emails, table_summaries_Emails, texts_Emails, table_list_Emails, img_base64_list_Emails, image_summaries_Emails = extraction_pipeline(text_list_Emails, table_list_Emails, save_img_folder)
                retriever_multi_vector_img = create_multi_vector_retriever(vectordb, text_summaries_Emails,texts_Emails,
                                                                    table_summaries_Emails,table_list_Emails,
                                                                    image_summaries_Emails,img_base64_list_Emails,)

        

        if retriever_multi_vector_img:
            main(retriever_multi_vector_img, username,guidelines_options,selected_guidelines, selected_document)
    else:
        st.title("Just Create")
        st.write("Welcome to Just Create!")
        st.write("With this AI powered solution you can generate and validate content according to brand guidelines in simple steps.")
        st.write("You can generate textual content just by asking with support of text and images.")
        st.write("If you want to validate content, you can provide text or image to validate or upload th PDF document itself.")

